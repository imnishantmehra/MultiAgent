# tools.py
from typing import Dict, Optional
import pandas as pd
import speech_recognition as sr
# from pydx2 import PdfReader
import docx2txt
import markdown
import re
from moviepy import VideoFileClip
from openpyxl import load_workbook
import json
from pptx import Presentation
import os
from crewai_tools import PDFSearchTool, Tool
from PyPDF2 import PdfReader
from pathlib import Path
from datetime import datetime



PLATFORM_LIMITS = {
    "twitter": {"chars": None, "words": 280},
    "instagram": {"chars": None, "words": 400},
    "linkedin": {"chars": None, "words": 600},
    "facebook": {"chars": None, "words": 1000},
    "wordpress": {"chars": None, "words": 2000},
    "youtube": {"chars": None, "words": 2000},
    "tiktok": {"chars": None, "words": 400},
}

UPLOAD_DIR = './uploads'

# Helper function to dynamically update the PDF path
def create_pdf_tool(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File {file_path} does not exist.")
    return PDFSearchTool(pdf=file_path)





class FileProcessor:
    def __init__(self):
        self.supported_formats = {
            # Document formats
            '.pdf': self.extract_from_pdf,
            '.docx': self.extract_from_docx,
            '.txt': self.extract_from_txt,
            '.md': self.extract_from_markdown,
            
            # Spreadsheet formats
            '.xlsx': self.extract_from_excel,
            '.xls': self.extract_from_excel,
            '.csv': self.extract_from_csv,
            
            # Presentation formats
            '.pptx': self.extract_from_powerpoint,
            '.ppt': self.extract_from_powerpoint,
            
            # Audio formats
            '.mp3': self.extract_from_audio,
            '.wav': self.extract_from_audio,
            '.m4a': self.extract_from_audio,
            
            # Video formats
            '.mp4': self.extract_from_video,
            '.mov': self.extract_from_video,
            '.avi': self.extract_from_video,
            
            # Web formats
            '.json': self.extract_from_json,
            '.html': self.extract_from_html,
        }

    def extract_text_from_file(self, file_path: str) -> str:
        """Main function to extract text from various file formats"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_extension = Path(file_path).suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        try:
            return self.supported_formats[file_extension](file_path)
        except Exception as e:
            raise Exception(f"Error extracting text from {file_path}: {str(e)}")

    def extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            reader = PdfReader(file_path)
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"PDF extraction error: {str(e)}")

    def extract_from_excel(self, file_path: str) -> str:
        """Extract text from Excel file"""
        try:
            # Read all sheets
            excel_data = []
            df = pd.read_excel(file_path, sheet_name=None)
            
            for sheet_name, sheet_data in df.items():
                excel_data.append(f"\nSheet: {sheet_name}")
                
                # Convert headers and data to string
                headers = sheet_data.columns.tolist()
                excel_data.append(f"Headers: {', '.join(str(h) for h in headers)}")
                
                # Convert data to string format
                for _, row in sheet_data.iterrows():
                    row_data = [str(cell) for cell in row]
                    excel_data.append(" | ".join(row_data))
            
            return "\n".join(excel_data)
        except Exception as e:
            raise Exception(f"Excel extraction error: {str(e)}")

    def extract_from_csv(self, file_path: str) -> str:
        """Extract text from CSV file"""
        try:
            df = pd.read_csv(file_path)
            csv_data = []
            
            # Convert headers and data to string
            headers = df.columns.tolist()
            csv_data.append(f"Headers: {', '.join(str(h) for h in headers)}")
            
            # Convert data to string format
            for _, row in df.iterrows():
                row_data = [str(cell) for cell in row]
                csv_data.append(" | ".join(row_data))
            
            return "\n".join(csv_data)
        except Exception as e:
            raise Exception(f"CSV extraction error: {str(e)}")

    def extract_from_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text_runs = []
            
            for i, slide in enumerate(prs.slides):
                text_runs.append(f"\nSlide {i+1}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            
            return "\n".join(text_runs)
        except Exception as e:
            raise Exception(f"PowerPoint extraction error: {str(e)}")

    def extract_from_audio(self, file_path: str) -> str:
        """Extract text from audio file using speech recognition"""
        try:
            # Initialize recognizer
            recognizer = sr.Recognizer()
            
            # Convert audio to WAV if it's not already
            if not file_path.lower().endswith('.wav'):
                from pydub import AudioSegment
                audio = AudioSegment.from_file(file_path)
                wav_path = file_path + '.wav'
                audio.export(wav_path, format='wav')
                file_path = wav_path
            
            # Perform speech recognition
            with sr.AudioFile(file_path) as source:
                audio = recognizer.record(source)
                text = recognizer.recognize_google(audio)
                
            # Clean up temporary WAV file if created
            if file_path.endswith('.wav') and os.path.exists(file_path):
                os.remove(file_path)
                
            return text
        except Exception as e:
            raise Exception(f"Audio extraction error: {str(e)}")

    def extract_from_video(self, file_path: str) -> str:
        """Extract text from video file by converting audio and using speech recognition"""
        try:
            # Extract audio from video
            video = VideoFileClip(file_path)
            audio = video.audio
            
            # Save audio temporarily
            temp_audio = file_path + '.wav'
            audio.write_audiofile(temp_audio)
            
            # Extract text from audio
            text = self.extract_from_audio(temp_audio)
            
            # Cleanup
            video.close()
            audio.close()
            if os.path.exists(temp_audio):
                os.remove(temp_audio)
                
            return text
        except Exception as e:
            raise Exception(f"Video extraction error: {str(e)}")

    def extract_from_json(self, file_path: str) -> str:
        """Extract text from JSON file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                return json.dumps(data, indent=2)
        except Exception as e:
            raise Exception(f"JSON extraction error: {str(e)}")

    def extract_from_html(self, file_path: str) -> str:
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                return soup.get_text(separator='\n', strip=True)
        except Exception as e:
            raise Exception(f"HTML extraction error: {str(e)}")

    def extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            text = docx2txt.process(file_path)
            return text
        except Exception as e:
            raise Exception(f"DOCX extraction error: {str(e)}")

    def extract_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            raise Exception(f"TXT extraction error: {str(e)}")

    def extract_from_markdown(self, file_path: str) -> str:
        """Extract text from Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
                html_content = markdown.markdown(md_content)
                return re.sub('<[^<]+?>', '', html_content)
        except Exception as e:
            raise Exception(f"Markdown extraction error: {str(e)}")

def generate_unique_content(content: str, week: int, day: str, platform: str) -> str:
    """
    Generate unique content based on week and day context
    """
    # Split content into sentences
    sentences = content.split('. ')
    
    # Use week and day to select different portions of content
    start_idx = ((week - 1) * 5 + ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"].index(day)) % len(sentences)
    
    # Take a different subset of sentences for each day
    num_sentences = min(5, len(sentences))  # Take up to 5 sentences
    selected_sentences = sentences[start_idx:start_idx + num_sentences]
    
    # If we need more sentences, wrap around to the beginning
    if len(selected_sentences) < num_sentences:
        selected_sentences.extend(sentences[:num_sentences - len(selected_sentences)])
    
    # Join sentences back together
    unique_content = '. '.join(selected_sentences)
    
    # Add platform-specific formatting
    if platform == "twitter":
        unique_content = f"Day  {unique_content}"
    elif platform == "instagram":
        unique_content = f"📱{unique_content}"
    elif platform == "linkedin":
        unique_content = f"💡 Professional Insight - {unique_content}"
    elif platform == "facebook":
        unique_content = f"🎯 Thought of the Day - {unique_content}"
    
    return unique_content


def generate_different_content(content, week, day, platform, post_number):
    """
    Generate unique content based on week, day context, and post number to ensure
    different content for multiple posts on the same day
    """
    # Split content into sentences
    sentences = content.split('. ')
    
    # Calculate different starting points for different posts on the same day
    base_idx = ((week - 1) * 5 + ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"].index(day)) % len(sentences)
    post_offset = (post_number - 1) * 3  # Offset by 3 sentences for each post
    start_idx = (base_idx + post_offset) % len(sentences)
    
    # Take a different subset of sentences for each post
    num_sentences = min(15, len(sentences))  # Take up to 5 sentences
    selected_sentences = sentences[start_idx:start_idx + num_sentences]
    
    # If we need more sentences, wrap around to the beginning
    if len(selected_sentences) < num_sentences:
        remaining = num_sentences - len(selected_sentences)
        selected_sentences.extend(sentences[:remaining])
    
    # Join sentences back together
    unique_content = '. '.join(selected_sentences)
    
    # Add platform-specific formatting with post number
    if platform == "twitter":
        unique_content = f"Update # - {unique_content}"
    elif platform == "instagram":
        unique_content = f"📱 \n{unique_content}"
    elif platform == "linkedin":
        unique_content = f"💡 Professional Insight # - {unique_content}"
    elif platform == "facebook":
        unique_content = f"🎯 Update # - {unique_content}"
    elif platform == "youtube":
        unique_content == f" {unique_content}"
    elif platform == "tiktok":
        unique_content == f" {unique_content}"
    
    return unique_content





def extract_title_from_content(content: str) -> str:
    """
    Extract or generate a title from the content
    """
    # Try to get the first sentence
    sentences = content.split('. ')
    if sentences:
        # Use the first sentence if it's not too long
        first_sentence = sentences[0].strip()
        if len(first_sentence) <= 50:
            return first_sentence
        
        # If first sentence is too long, take first few words
        words = first_sentence.split()
        title = ' '.join(words[:7])
        return f"{title}..."
    
    return "Untitled Post"

def process_content_for_platform(content: str, platform: str, limits: Dict) -> str:
    """
    Process and format content according to platform-specific limits
    """
    # Clean content
    content = content.strip()
    
    # Apply character limit if specified
    if limits.get("chars"):
        content = content[:limits["chars"]]
    
    # Apply word limit if specified
    if limits.get("words"):
        words = content.split()
        content = ' '.join(words[:limits["words"]])
    
    # Add platform-specific formatting
    if platform == "twitter":
        # Leave room for hashtags
        if len(content) > 240:
            content = content[:237] + "..."
        content += "#Content #Social"
    
    elif platform == "instagram":
        content += "#Instagram #Social"
    
    elif platform == "linkedin":
        content += "#Professional #Development"
    
    elif platform == "facebook":
        content += "Like and share if you agree! 👍"

    elif platform == "youtube":
        content += "#Professional #Development"

    elif platform == "tiktok":
        content == "follow for more"
    
    return content.strip()


# Instantiate FileProcessor
file_processor = FileProcessor()

# Define a tool wrapper for file processing
file_extraction_tool = Tool(
    name="FileExtractionTool",
    description="A universal tool for extracting text from various file formats.",
    func=lambda file_path: file_processor.extract_text_from_file(file_path),  # Use 'func' as the field name
)



def trim_content(content: str, platform: str, limits: dict) -> str:
    """Trim content based on platform-specific character or word limits."""
    char_limit = limits.get(platform, {}).get("chars")
    word_limit = limits.get(platform, {}).get("words")
    
    # Apply character limit
    if char_limit:
        content = content[:char_limit].strip()
    
    # Apply word limit
    if word_limit:
        words = content.split()
        content = " ".join(words[:word_limit]).strip()
    
    return content



def generate_twitter_post(content: str, limits: dict) -> str:
    """Generate a concise Twitter post with dynamic hashtags."""
    content = trim_content(content, "twitter", limits)
    hashtags = "#Trending #Innovation #Tech"
    return f"{content}\n\n{hashtags}"

twitter_post_tool = Tool(
    name="TwitterPostGeneratorTool",
    description="Generates concise and impactful tweets with dynamic hashtags, respecting platform limits.",
    func=lambda content: generate_twitter_post(content, PLATFORM_LIMITS),
)



def generate_instagram_post(content: str, limits: dict) -> str:
    """Generate an Instagram post with emojis and hashtags."""
    content = trim_content(content, "instagram", limits)
    hashtags = "#InstaLife #Photography #Explore"
    emoji = "📸✨"
    return f"{emoji} {content.strip()} {hashtags}"

instagram_post_tool = Tool(
    name="InstagramPostGeneratorTool",
    description="Generates Instagram captions with emojis and hashtags, respecting platform limits.",
    func=lambda content: generate_instagram_post(content, PLATFORM_LIMITS),
)


def generate_linkedin_post(content: str, limits: dict) -> str:
    """Generate a professional LinkedIn post with hashtags."""
    content = trim_content(content, "linkedin", limits)
    hashtags = "#Professional #Leadership #Growth"
    return f"{content.strip()} {hashtags} Engage with this post by sharing your thoughts!"

linkedin_post_tool = Tool(
    name="LinkedInPostGeneratorTool",
    description="Generates professional LinkedIn posts with dynamic formatting and hashtags.",
    func=lambda content: generate_linkedin_post(content, PLATFORM_LIMITS),
)


def generate_facebook_post(content: str, limits: dict) -> str:
    """Generate a friendly Facebook post."""
    content = trim_content(content, "facebook", limits)
    return f"{content.strip()} Like and share if you agree! 👍"

facebook_post_tool = Tool(
    name="FacebookPostGeneratorTool",
    description="Generates friendly and engaging Facebook posts, respecting platform limits.",
    func=lambda content: generate_facebook_post(content, PLATFORM_LIMITS),
)


def generate_wordpress_post(content: str, limits: dict) -> str:
    """Generate a WordPress blog post with SEO-friendly structure."""
    content = trim_content(content, "wordpress", limits)
    # title = "Blog Title: " + (content[:50] + "..." if len(content) > 50 else content)
    return f" {content.strip()} #SEO #WordPress"

wordpress_post_tool = Tool(
    name="WordPressPostGeneratorTool",
    description="Generates SEO-optimized WordPress blog posts, respecting platform limits.",
    func=lambda content: generate_wordpress_post(content, PLATFORM_LIMITS),
)


def generate_youtube_post(content: str, limits: dict) -> str:
    """Generate a youtube post which is further used for image or video genration"""
    content = trim_content(content, "youtube", limits)
    # title = "Blog Title: " + (content[:50] + "..." if len(content) > 50 else content)
    return f" {content.strip()} "

youtube_post_tool = Tool(
    name="youtubePostGeneratorTool",
    description="Generate a youtube post which is further used for image or video genration.",
    func=lambda content: generate_youtube_post(content, PLATFORM_LIMITS),
)


def generate_tiktok_post(content: str, limits: dict) -> str:
    """Generate a tiktok post which is further used for image or video genration"""
    content = trim_content(content, "tiktok", limits)
    # title = "Blog Title: " + (content[:50] + "..." if len(content) > 50 else content)
    return f" {content.strip()} "

tiktok_post_tool = Tool(
    name="tiktokPostGeneratorTool",
    description="Generate a youtube post which is further used for image or video genration.",
    func=lambda content: generate_tiktok_post(content, PLATFORM_LIMITS),
)
